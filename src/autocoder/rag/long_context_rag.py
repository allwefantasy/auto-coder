import json
import os
from concurrent.futures import ThreadPoolExecutor, as_completed
from io import BytesIO
from typing import Any, Dict, Generator, List, Optional, Tuple

import byzerllm
import docx2txt
import pandas as pd
import pathspec
import pptx
from byzerllm import ByzerLLM
from docx import Document
from jinja2 import Template
from loguru import logger
from rich.console import Console
from rich.table import Table
from rich.panel import Panel
from rich.text import Text
from openai import OpenAI
from openpyxl import load_workbook
from pypdf import PdfReader
import time
from byzerllm.utils.client.code_utils import extract_code

from autocoder.common import AutoCoderArgs, SourceCode


class LongContextRAG:
    def __init__(self, llm: ByzerLLM, args: AutoCoderArgs, path: str) -> None:
        self.llm = llm
        self.args = args
        self.path = path

        self.tokenizer = None
        if llm.is_model_exist("deepseek_tokenizer"):
            self.tokenizer = ByzerLLM()
            self.tokenizer.setup_default_model_name("deepseek_tokenizer")

        self.required_exts = (
            [ext.strip() for ext in self.args.required_exts.split(",")]
            if self.args.required_exts
            else []
        )

        if args.rag_url and args.rag_url.startswith("http://"):
            if not args.rag_token:
                raise ValueError(
                    "You are in client mode, please provide the RAG token. e.g. rag_token: your_token_here"
                )
            self.client = OpenAI(api_key=args.rag_token, base_url=args.rag_url)
        else:
            self.client = None
            # if not pure client mode, then the path should be provided
            if (
                not self.path
                and args.rag_url
                and not args.rag_url.startswith("http://")
            ):
                self.path = args.rag_url

            if not self.path:
                raise ValueError(
                    "Please provide the path to the documents in the local file system."
                )

        self.ignore_spec = self._load_ignore_file()

        self.token_limit = 120000

        ## 检查当前目录下所有文件是否超过 120k tokens ，并且打印出来
        self.token_exceed_files = []
        if self.tokenizer is not None:
            docs = self._retrieve_documents()
            for doc in docs:
                token_num = self.count_tokens(doc.source_code)
                if token_num > self.token_limit:
                    self.token_exceed_files.append(doc.module_name)

        if self.token_exceed_files:
            logger.warning(
                f"以下文件超过了 120k tokens: {self.token_exceed_files},将无法使用 RAG 模型进行搜索。"
            )

    def count_tokens(self, text: str) -> int:
        if self.tokenizer is None:
            return -1
        try:
            v = self.tokenizer.chat_oai(
                conversations=[{"role": "user", "content": text}]
            )
            return int(v[0].output)
        except Exception as e:
            logger.error(f"Error counting tokens: {str(e)}")
            return -1

    def extract_text_from_pdf(self, pdf_content):
        pdf_file = BytesIO(pdf_content)
        pdf_reader = PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text

    def extract_text_from_docx(self, docx_content):
        docx_file = BytesIO(docx_content)
        text = docx2txt.process(docx_file)
        return text

    def extract_text_from_excel(self, excel_path) -> List[Tuple[str, str]]:
        sheet_list = []
        wb = load_workbook(excel_path)
        tmpl = Template(
            """{% for row in rows %}
{% for cell in row %}"{{ cell }}"{% if not loop.last %},{% endif %}{% endfor %}{% if not loop.last %}{% endif %}{% endfor %}
        """
        )
        for ws in wb:
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            # 过滤掉rows中全是null的行
            rows = [row for row in rows if any(row)]
            # 所有的None都转换成空字符串
            rows = [[cell if cell is not None else "" for cell in row] for row in rows]
            content = tmpl.render(rows=rows)
            sheet_list.append([excel_path + f"#{ws.title}", content])
        return sheet_list

    def extract_text_from_ppt(self, ppt_path) -> List[Tuple[str, str]]:
        presentation = pptx.Presentation(ppt_path)
        text_template = Template("""文本:{{ paragraphs }}""")
        table_template = Template(
            """表单:{% for row in rows %}
{% for cell in row %}"{{ cell }}"{% if not loop.last %},{% endif %}{% endfor %}{% if not loop.last %}{% endif %}{% endfor %}"""
        )

        slide_list = []
        for slide in presentation.slides:
            shape_list = []
            for shape in slide.shapes:
                contents = []
                # TODO: support charts, images
                # support text
                if shape.has_text_frame:
                    paragraphs = []
                    for paragraph in shape.text_frame.paragraphs:
                        paragraphs.append(paragraph.text)
                    if len("".join(paragraphs).strip()) > 0:
                        contents.append("".join(paragraphs))
                if len(contents) > 0:
                    shape_list.append(
                        text_template.render(paragraphs="\n".join(contents))
                    )
                # support table
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [
                            (cell.text_frame.text).strip() for cell in row.cells
                        ]
                        table_data.append(row_data)
                    shape_list.append(table_template.render(rows=table_data))
            if len(shape_list) > 0:
                slide_list.append(
                    [ppt_path + f"#{slide.slide_id}", "\n".join(shape_list)]
                )
        return slide_list

    @byzerllm.prompt()
    def _check_relevance(self, query: str, documents: List[str]) -> str:
        """
        请判断以下文档是否能够回答给出的问题。

        文档：
        {% for doc in documents %}
        {{ doc }}
        {% endfor %}

        问题：{{ query }}

        如果该文档提供的知识能够回答问题，那么请回复"yes" 否则回复"no"。
        """

    @byzerllm.prompt()
    def _check_relevance_with_conversation(
        self, conversations: List[Dict[str, str]], documents: List[str]
    ) -> str:
        """
        使用以下文档和对话历史来回答问题。如果文档中没有相关信息，请说"我没有足够的信息来回答这个问题"。

        文档：
        {% for doc in documents %}
        {{ doc }}
        {% endfor %}

        对话历史：
        {% for msg in conversations %}
        <{{ msg.role }}>: {{ msg.content }}
        {% endfor %}

        请结合提供的文档以及用户对话历史，判断提供的文档是不是能回答用户的最后一个问题。
        如果该文档提供的知识能够回答问题，那么请回复"yes" 否则回复"no"。
        """

    @byzerllm.prompt()
    def extract_relevance_info_from_docs_with_conversation(
        self, conversations: List[Dict[str, str]], documents: List[str]
    ) -> str:
        """
        使用以下文档和对话历史来提取相关信息。

        文档：
        {% for doc in documents %}
        {{ doc }}
        {% endfor %}

        对话历史：
        {% for msg in conversations %}
        <{{ msg.role }}>: {{ msg.content }}
        {% endfor %}

        请根据提供的文档内容、用户对话历史以及最后一个问题，提取并总结文档中与问题相关的重要信息。
        如果文档中没有相关信息，请回复"该文档中没有与问题相关的信息"。
        提取的信息尽量保持和原文中的一样，并且只输出这些信息。
        """
    
    @byzerllm.prompt()
    def extract_relevance_range_from_docs_with_conversation(
        self, conversations: List[Dict[str, str]], documents: List[str]
    ) -> str:
        """
        根据提供的文档和对话历史提取相关信息范围。

        输入:
        1. 文档内容:
        {% for doc in documents %}
        {{ doc }}
        {% endfor %}

        2. 对话历史:
        {% for msg in conversations %}
        <{{ msg.role }}>: {{ msg.content }}
        {% endfor %}

        任务:
        1. 分析最后一个用户问题及其上下文。
        2. 在文档中找出与问题相关的一个或多个重要信息段。
        3. 对每个相关信息段，确定其起始行号(start_line)和结束行号(end_line)。
        4. 信息段数量不超过4个。

        输出要求:
        1. 返回一个JSON数组，每个元素包含"start_line"和"end_line"。
        2. start_line和end_line必须是整数，表示文档中的行号。
        3. 行号从1开始计数。
        4. 如果没有相关信息，返回空数组[]。

        输出格式:
        严格的JSON数组，不包含其他文字或解释。

        示例:
        1.  文档：
            1 这是这篇动物科普文。
            2 大象是陆地上最大的动物之一。
            3 它们生活在非洲和亚洲。
            问题：大象生活在哪里？
            返回：[{"start_line": 2, "end_line": 3}]

        2.  文档：
            1 地球是太阳系第三行星，
            2 有海洋、沙漠，温度适宜，
            3 是已知唯一有生命的星球。
            4 太阳则是太阳系的唯一恒心。
            问题：地球的特点是什么？
            返回：[{"start_line": 1, "end_line": 3}]

        3.  文档：
            1 苹果富含维生素。
            2 香蕉含有大量钾元素。
            问题：橙子的特点是什么？
            返回：[]        
        """

    @byzerllm.prompt()
    def _answer_question(
        self, query: str, relevant_docs: List[str]
    ) -> Generator[str, None, None]:
        """
        使用以下文档来回答问题。如果文档中没有相关信息，请说"我没有足够的信息来回答这个问题"。

        文档：
        {% for doc in relevant_docs %}
        {{ doc }}
        {% endfor %}

        问题：{{ query }}

        回答：
        """

    def _load_ignore_file(self):
        serveignore_path = os.path.join(self.path, ".serveignore")
        gitignore_path = os.path.join(self.path, ".gitignore")

        if os.path.exists(serveignore_path):
            with open(serveignore_path, "r") as ignore_file:
                return pathspec.PathSpec.from_lines("gitwildmatch", ignore_file)
        elif os.path.exists(gitignore_path):
            with open(gitignore_path, "r") as ignore_file:
                return pathspec.PathSpec.from_lines("gitwildmatch", ignore_file)
        return None

    def _retrieve_documents(self) -> List[SourceCode]:
        documents = []
        for root, dirs, files in os.walk(self.path):
            # 过滤掉隐藏目录
            dirs[:] = [d for d in dirs if not d.startswith(".")]

            # 应用 .serveignore 或 .gitignore 规则
            if self.ignore_spec:
                relative_root = os.path.relpath(root, self.path)
                dirs[:] = [
                    d
                    for d in dirs
                    if not self.ignore_spec.match_file(os.path.join(relative_root, d))
                ]
                files = [
                    f
                    for f in files
                    if not self.ignore_spec.match_file(os.path.join(relative_root, f))
                ]

            for file in files:
                if self.required_exts:
                    if not any(file.endswith(ext) for ext in self.required_exts):
                        continue

                file_path = os.path.join(root, file)
                relative_path = os.path.relpath(file_path, self.path)

                try:
                    if file.endswith(".pdf"):
                        with open(file_path, "rb") as f:
                            content = self.extract_text_from_pdf(f.read())
                        documents.append(
                            SourceCode(module_name=file_path, source_code=content)
                        )
                    elif file.endswith(".docx"):
                        with open(file_path, "rb") as f:
                            content = self.extract_text_from_docx(f.read())
                        documents.append(
                            SourceCode(
                                module_name=f"##File: {file_path}", source_code=content
                            )
                        )
                    elif file.endswith(".xlsx") or file.endswith(".xls"):
                        sheets = self.extract_text_from_excel(file_path)
                        for sheet in sheets:
                            documents.append(
                                SourceCode(
                                    module_name=f"##File: {file_path}#{sheet[0]}",
                                    source_code=sheet[1],
                                )
                            )
                    elif file.endswith(".pptx"):
                        slides = self.extract_text_from_ppt(file_path)
                        content = ""
                        for slide in slides:
                            content += f"#{slide[0]}\n{slide[1]}\n\n"
                        documents.append(
                            SourceCode(
                                module_name=f"##File: {file_path}", source_code=content
                            )
                        )
                    else:
                        with open(file_path, "r", encoding="utf-8") as f:
                            content = f.read()
                            documents.append(
                                SourceCode(
                                    module_name=f"##File: {file_path}",
                                    source_code=content,
                                )
                            )
                except Exception as e:
                    logger.error(f"Error processing file {file_path}: {str(e)}")

        return documents

    def build(self):
        pass

    def search(self, query: str) -> List[SourceCode]:
        target_query = query
        only_contexts = False
        if self.args.enable_rag_search and isinstance(self.args.enable_rag_search, str):
            target_query = self.args.enable_rag_search
        elif self.args.enable_rag_context and isinstance(
            self.args.enable_rag_context, str
        ):
            target_query = self.args.enable_rag_context
            only_contexts = True

        logger.info("Search from RAG.....")
        logger.info(f"Query: {target_query} only_contexts: {only_contexts}")

        if self.client:
            new_query = json.dumps(
                {"query": target_query, "only_contexts": only_contexts},
                ensure_ascii=False,
            )
            response = self.client.chat.completions.create(
                messages=[{"role": "user", "content": new_query}],
                model=self.args.model,
            )
            v = response.choices[0].message.content
            if not only_contexts:
                return [SourceCode(module_name=f"RAG:{target_query}", source_code=v)]

            json_lines = [json.loads(line) for line in v.split("\n") if line.strip()]
            return [SourceCode.model_validate(json_line) for json_line in json_lines]
        else:
            if only_contexts:
                return self._filter_docs(target_query)
            else:
                v, contexts = self.stream_chat_oai(
                    conversations=[{"role": "user", "content": new_query}]
                )
                url = ",".join(contexts)
                return [SourceCode(module_name=f"RAG:{url}", source_code="".join(v))]

    def _filter_docs(self, conversations: List[Dict[str, str]]) -> List[SourceCode]:
        documents = self._retrieve_documents()
        with ThreadPoolExecutor(
            max_workers=self.args.index_filter_workers or 5
        ) as executor:
            future_to_doc = {}
            for doc in documents:
                content = self._check_relevance_with_conversation.prompt(
                    conversations, [doc.source_code]
                )
                if self.tokenizer and self.count_tokens(content) > self.token_limit:
                    logger.warning(
                        f"{doc.module_name} 以及对话上线文 文件超过 120k tokens，将无法使用 RAG 模型进行搜索。"
                    )
                    continue

                m = executor.submit(
                    self._check_relevance_with_conversation.with_llm(self.llm).run,
                    conversations,
                    [f"##File: {doc.module_name}\n{doc.source_code}"],
                )
                future_to_doc[m] = doc
        relevant_docs = []
        for future in as_completed(future_to_doc):
            try:
                doc = future_to_doc[future]
                v = future.result()
                if "yes" in v.strip().lower():
                    relevant_docs.append(doc)
            except Exception as exc:
                logger.error(f"Document processing generated an exception: {exc}")

        return relevant_docs
    
    def stream_chat_oai(
        self,
        conversations,
        model: Optional[str] = None,
        role_mapping=None,
        llm_config: Dict[str, Any] = {},
    ):
        try:
            return self._stream_chat_oai(
                conversations, model=model, role_mapping=role_mapping, llm_config=llm_config
            )
        except Exception as e:
            logger.error(f"Error in stream_chat_oai: {str(e)}")
            e.print_exc()
            return ["出现错误，请稍后再试。"], []

    def _stream_chat_oai(
        self,
        conversations,
        model: Optional[str] = None,
        role_mapping=None,
        llm_config: Dict[str, Any] = {},
    ):
        if self.client:
            model = model or self.args.model
            response = self.client.chat.completions.create(
                model=model,
                messages=conversations,
                stream=True,
            )

            def response_generator():
                for chunk in response:
                    if chunk.choices[0].delta.content is not None:
                        yield chunk.choices[0].delta.content

            return response_generator(), []
        else:
            query = conversations[-1]["content"]
            context = []

            if (
                "使用四到五个字直接返回这句话的简要主题，不要解释、不要标点、不要语气词、不要多余文本，不要加粗，如果没有主题"
                in query
                or "简要总结一下对话内容，用作后续的上下文提示 prompt，控制在 200 字以内"
                in query
            ):
                chunks = self.llm.stream_chat_oai(
                    conversations=conversations,
                    model=model,
                    role_mapping=role_mapping,
                    llm_config=llm_config,
                    delta_mode=True,
                )
                return (chunk[0] for chunk in chunks), context

            only_contexts = False
            try:
                v = json.loads(query)
                if "only_contexts" in v:
                    query = v["query"]
                    only_contexts = v["only_contexts"]
            except json.JSONDecodeError:
                pass

            logger.info(f"Query: {query} only_contexts: {only_contexts}")
            start_time = time.time()
            relevant_docs: List[SourceCode] = self._filter_docs(conversations)
            filter_time = time.time() - start_time

            logger.info(
                f"Filter time: {filter_time:.2f} seconds with {len(relevant_docs)} docs"
            )

            if only_contexts:
                return (doc.model_dump_json() + "\n" for doc in relevant_docs), []

            if not relevant_docs:
                return ["没有找到相关的文档来回答这个问题。"], []

            context = [doc.module_name for doc in relevant_docs]

            console = Console()

            # Create a table for the query information
            query_table = Table(title="Query Information", show_header=False)
            query_table.add_row("Query", query)
            query_table.add_row("Relevant docs", str(len(relevant_docs)))

            # Add relevant docs information
            relevant_docs_info = "\n".join([f"- {doc.module_name}" for doc in relevant_docs])
            query_table.add_row("Relevant docs list", relevant_docs_info)

            first_round_full_docs = []
            second_round_extracted_docs = []
            sencond_round_time = 0

            if self.tokenizer is not None:
                final_relevant_docs = []
                token_count = 0
                doc_num_count = 0
                for doc in relevant_docs:
                    doc_tokens = self.count_tokens(doc.source_code)
                    doc_num_count += 1
                    if token_count + doc_tokens <= self.token_limit:
                        final_relevant_docs.append(doc)
                        token_count += doc_tokens
                    else:
                        break

                ## 如果relevant_docs 拼接后超过了 token_limit 限制，那么我们会重新获取文档。
                ## 首先遍历文档，不断添加到 final_relevant_docs 直到达到了 token_limit * 0.6 的限制
                ## 剩下的文档，我们会继续做遍历，但是会对每个文档做信息抽取，然后继续添加到 final_relevant_docs 中
                ## 直到 token 数量达到 token_limit * 0.4 的限制

                if len(final_relevant_docs) < len(relevant_docs):
                    ## 过滤出前面80% token的文档
                    token_count = 0
                    new_token_limit = self.token_limit * 0.8
                    doc_num_count = 0
                    for doc in relevant_docs:
                        doc_tokens = self.count_tokens(doc.source_code)
                        doc_num_count += 1
                        if token_count + doc_tokens <= new_token_limit:
                            first_round_full_docs.append(doc)
                            token_count += doc_tokens
                        else:
                            break

                    ## 获取剩下的token数量 和 文档
                    sencond_round_start_time = time.time()
                    remaining_tokens = self.token_limit - new_token_limit
                    remaining_docs = relevant_docs[len(first_round_full_docs) :]                    

                    with ThreadPoolExecutor(
                        max_workers=self.args.index_filter_workers or 5
                    ) as executor:                        
                        def process_doc(doc):                            
                            extracted_info = self.extract_relevance_info_from_docs_with_conversation.with_llm(
                                self.llm
                            ).run(
                                conversations, [doc.source_code]
                            )                            
                            
                            return SourceCode(
                                module_name=doc.module_name, source_code=extracted_info
                            )
                        
                        def process_range_doc(doc, max_retries=3):
                            for attempt in range(max_retries):
                                content = ""
                                try:
                                    source_code_with_line_number = ""          
                                    source_code_lines = doc.source_code.split("\n")  
                                    for idx, line in enumerate(source_code_lines):
                                        source_code_with_line_number += f"{idx+1} {line}\n"                

                                    extracted_info = self.extract_relevance_range_from_docs_with_conversation.with_llm(
                                        self.llm
                                    ).run(
                                        conversations, [source_code_with_line_number]
                                    )
                                                                                            
                                    json_str = extract_code(extracted_info)[0][1]
                                    json_objs = json.loads(json_str)                                    
                                                                    
                                    for json_obj in json_objs:
                                        start_line = json_obj["start_line"] - 1
                                        end_line = json_obj["end_line"]
                                        chunk = "\n".join(source_code_lines[start_line:end_line])
                                        content += chunk + "\n"
                                    
                                    return SourceCode(
                                        module_name=doc.module_name, source_code=content.strip()
                                    )
                                except Exception as e:
                                    if attempt < max_retries - 1:
                                        logger.warning(f"Error processing doc {doc.module_name}, retrying... (Attempt {attempt + 1}) Error: {str(e)}")                                        
                                    else:
                                        logger.error(f"Failed to process doc {doc.module_name} after {max_retries} attempts: {str(e)}")
                                        return SourceCode(
                                            module_name=doc.module_name, source_code=content.strip()
                                        )

                        future_to_doc = {}
                        for doc in remaining_docs:
                            future = executor.submit(process_range_doc, doc)
                            future_to_doc[future] = doc

                        for future in as_completed(future_to_doc):
                            doc = future_to_doc[future]
                            try:
                                result = future.result()
                                if result and remaining_tokens > 0:
                                    second_round_extracted_docs.append(result)
                                    tokens = self.count_tokens(result.source_code)
                                    if tokens > 0:
                                        remaining_tokens -= tokens
                                    else:
                                        logger.warning(f"Token count for doc {doc.module_name} is 0 or negative")
                            except Exception as exc:
                                logger.error(
                                    f"Processing doc {doc.module_name} generated an exception: {exc}"
                                )
                    final_relevant_docs = (
                        first_round_full_docs + second_round_extracted_docs
                    )
                    sencond_round_time = time.time() - sencond_round_start_time
                relevant_docs = final_relevant_docs
            else:
                relevant_docs = relevant_docs[: self.args.index_filter_file_num]

            logger.info(f"Finally send to model: {len(relevant_docs)}")

            query_table.add_row("Only contexts", str(only_contexts))
            query_table.add_row("Filter time", f"{filter_time:.2f} seconds")
            query_table.add_row("Final relevant docs", str(len(relevant_docs)))
            query_table.add_row("first_round_full_docs", str(len(first_round_full_docs)))
            query_table.add_row("second_round_extracted_docs", str(len(second_round_extracted_docs)))
            query_table.add_row("Second round time", f"{sencond_round_time:.2f} seconds")

            # Add relevant docs information
            final_relevant_docs_info = "\n".join([f"- {doc.module_name}" for doc in relevant_docs])
            query_table.add_row("Final Relevant docs list", final_relevant_docs_info)

            # Create a panel to contain the table
            panel = Panel(
                query_table,
                title="RAG Search Results",
                expand=False,
            )

            # Log the panel using rich
            console.print(panel)

            logger.info(f"Start to send to model {model}")

            new_conversations = conversations[:-1] + [
                {
                    "role": "user",
                    "content": self._answer_question.prompt(
                        query=query,
                        relevant_docs=[doc.source_code for doc in relevant_docs],
                    ),
                }
            ]

            # # 将 new_conversations 转化为 JSON 并写入文件
            # with open('/tmp/rag.json', 'w', encoding='utf-8') as f:
            #     json.dump(new_conversations, f, ensure_ascii=False, indent=2)

            chunks = self.llm.stream_chat_oai(
                conversations=new_conversations,
                model=model,
                role_mapping=role_mapping,
                llm_config=llm_config,
                delta_mode=True,
            )
            return (chunk[0] for chunk in chunks), context
