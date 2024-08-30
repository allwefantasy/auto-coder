from typing import List, Tuple
import pptx
from jinja2 import Template

def extract_text_from_ppt(ppt_path) -> List[Tuple[str, str]]:
    presentation = pptx.Presentation(ppt_path)
    text_template = Template("""文本:{{ paragraphs }}""")
    table_template = Template(
        """表单:{% for row in rows %}
{% for cell in row %}"{{ cell }}"{% if not loop.last %},{% endif %}{% endfor %}{% if not loop.last %}{% endif %}{% endfor %}"""
    )

    slide_list = []
    for slide in presentation.slides:
        shape_list = []
        for shape in slide.shapes:
            contents = []
            # TODO: support charts, images
            # support text
            if shape.has_text_frame:
                paragraphs = []
                for paragraph in shape.text_frame.paragraphs:
                    paragraphs.append(paragraph.text)
                if len("".join(paragraphs).strip()) > 0:
                    contents.append("".join(paragraphs))
            if len(contents) > 0:
                shape_list.append(
                    text_template.render(paragraphs="\n".join(contents))
                )
            # support table
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [
                        (cell.text_frame.text).strip() for cell in row.cells
                    ]
                    table_data.append(row_data)
                shape_list.append(table_template.render(rows=table_data))
        if len(shape_list) > 0:
            slide_list.append(
                [ppt_path + f"#{slide.slide_id}", "\n".join(shape_list)]
            )
    return slide_list